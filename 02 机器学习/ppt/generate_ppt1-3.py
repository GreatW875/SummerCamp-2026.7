#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""机器学习学习汇报PPT生成脚本 | Ubuntu 22.04 + conda(sc) + python-pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY_D = RGBColor(0x33, 0x33, 0x33)
GRAY_M = RGBColor(0x66, 0x66, 0x66)
GRAY_L = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0, 0xB4, 0xD8)
ACCENT2 = RGBColor(0xFF, 0x6B, 0x6B)
TH_BG = RGBColor(0x1A, 0x1A, 0x1A)
TR_BG = RGBColor(0x0D, 0x0D, 0x0D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BL = prs.slide_layouts[6]


def set_bg(s, c=BLACK):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()


def _set_font(run, size, bold, color, font='微软雅黑'):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)


def add_text(s, l, t, w, h, text, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    _set_font(run, size, bold, color)
    return tb


def title_bar(s, title, sub=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.08), Inches(0.45))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    add_text(s, Inches(0.85), Inches(0.45), Inches(11), Inches(0.6), title, 28, True)
    if sub:
        add_text(s, Inches(0.85), Inches(0.95), Inches(11), Inches(0.35), sub, 13, False, GRAY_L)
    sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.1), Emu(9525))
    sep.fill.solid(); sep.fill.fore_color.rgb = GRAY_D; sep.line.fill.background()


def section_label(s, text, top=1.6):
    add_text(s, Inches(0.6), Inches(top), Inches(5), Inches(0.35), text, 12, True, ACCENT)


def _set_cell_border(cell, color='333333'):
    tcPr = cell._tc.get_or_add_tcPr()
    for bn in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = tcPr.find(qn(bn))
        if ln is None:
            ln = etree.SubElement(tcPr, qn(bn))
        ln.set('w', '6350')
        sf = ln.find(qn('a:solidFill'))
        if sf is None:
            sf = etree.SubElement(ln, qn('a:solidFill'))
        sc = sf.find(qn('a:srgbClr'))
        if sc is None:
            sc = etree.SubElement(sf, qn('a:srgbClr'))
        sc.set('val', color)


def add_table(s, l, t, w, h, data, cw=None, hs=12, cs=11, fcb=True):
    rows, cols = len(data), len(data[0])
    ts = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = ts.table
    if cw:
        for i, wi in enumerate(cw):
            tbl.columns[i].width = wi
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c); cell.text = ''
            tf = cell.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_right = Emu(45720)
            tf.margin_top = tf.margin_bottom = Emu(22860)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = str(data[r][c])
            is_h = (r == 0)
            is_fcb = fcb and c == 0 and r > 0
            sz = hs if is_h else cs
            bold = is_h or is_fcb
            clr = ACCENT if is_h else (ACCENT2 if is_fcb else WHITE)
            _set_font(run, sz, bold, clr)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TH_BG if is_h else (TR_BG if r % 2 == 1 else BLACK)
            _set_cell_border(cell)
    return tbl


def bullets(s, l, t, w, h, items, size=12, color=GRAY_L):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(5)
        run = p.add_run(); run.text = f'▸  {item}'
        _set_font(run, size, False, color)
    return tb


# ========== Slide 1: 封面 ==========
s = prs.slides.add_slide(BL); set_bg(s)
for y in [2.2, 2.25, 2.3]:
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(y), Inches(0.6), Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
add_text(s, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2), '机器学习学习汇报', 54, True)
add_text(s, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
         '基础原理 · 树模型与集成 · 特征工程', 22, False, GRAY_L)
add_text(s, Inches(1.5), Inches(6.2), Inches(6), Inches(0.4), '暑期培训 · 机器学习方向', 14, False, GRAY_M)
add_text(s, Inches(1.5), Inches(6.6), Inches(6), Inches(0.4), '2026', 12, False, GRAY_D)
vl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.5), Inches(2.5), Emu(9525), Inches(3))
vl.fill.solid(); vl.fill.fore_color.rgb = GRAY_D; vl.line.fill.background()

# ========== Slide 2: 目录 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '目录 / CONTENTS')
toc = [
    ('01', '基础与原理', '监督/无监督 · 线性模型 · 损失函数 · KNN/K-Means · SVM'),
    ('02', '树模型与集成学习', '决策树 · 随机森林 · GBDT/XGBoost · LightGBM · CatBoost · Stacking'),
    ('03', '特征工程与数据处理', 'EDA · 编码 · 变换 · 构造 · 选择 · 不平衡 · Pipeline'),
]
for i, (n, t, d) in enumerate(toc):
    y = 2.0 + i * 1.6
    add_text(s, Inches(1.0), Inches(y), Inches(1.2), Inches(0.8), n, 42, True, ACCENT)
    add_text(s, Inches(2.5), Inches(y + 0.05), Inches(8), Inches(0.5), t, 22, True)
    add_text(s, Inches(2.5), Inches(y + 0.55), Inches(9), Inches(0.4), d, 13, False, GRAY_L)
    if i < len(toc) - 1:
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(y + 1.3), Inches(11), Emu(6350))
        sp.fill.solid(); sp.fill.fore_color.rgb = GRAY_D; sp.line.fill.background()

# ========== Slide 3: 基本概念 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '阶段一：基础与原理', '机器学习 · 神经网络 · 深度学习')
section_label(s, '三者关系')
add_text(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(0.4),
         '机器学习 ⊃ 神经网络 ⊃ 深度学习', 16, True, ACCENT2)
d1 = [
    ['层级', '范围', '核心'],
    ['机器学习', '最外层', '从数据中学习规律的算法总称'],
    ['神经网络', '中间层', '模拟神经元连接的机器学习方法'],
    ['深度学习', '最内层', '层数多的神经网络（>2隐藏层）'],
]
add_table(s, Inches(0.6), Inches(2.5), Inches(5.5), Inches(2.2), d1,
          cw=[Inches(1.2), Inches(1.5), Inches(2.8)], cs=11)
section_label(s, '监督 vs 无监督', top=1.6)
d2 = [
    ['维度', '监督学习', '无监督学习'],
    ['数据', '有标签 (X+y)', '无标签 (只有X)'],
    ['目标', '学 X→y 映射', '发现内在结构'],
    ['任务', '分类、回归', '聚类、降维、异常检测'],
    ['代表', 'LR、SVM、RF', 'K-Means、PCA'],
]
add_table(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(2.7), d2,
          cw=[Inches(1.0), Inches(2.3), Inches(2.6)], cs=11)
section_label(s, '数据集划分', top=4.9)
d3 = [
    ['集合', '作用', '占比', '类比'],
    ['训练集', '学习参数(w,b)', '70~80%', '平时作业'],
    ['验证集', '调超参数、选模型', '从训练集划出', '模拟考试'],
    ['测试集', '最终评估泛化能力', '20~30%', '期末考试（仅用一次）'],
]
add_table(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.8), d3,
          cw=[Inches(1.2), Inches(3.5), Inches(2.0), Inches(5.4)], cs=11)

# ========== Slide 4: 参数/过拟合/正则化 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '核心概念：参数 · 过拟合 · 正则化')
section_label(s, '参数 vs 超参数')
d1 = [
    ['类型', '含义', '谁决定', '例子'],
    ['参数', '模型从数据中学到的', '模型自己', 'w、b、神经网络权重'],
    ['超参数', '训练前人为设定的', '人', '学习率、正则化系数、树深'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(1.8), d1,
          cw=[Inches(1.0), Inches(2.0), Inches(1.2), Inches(1.6)], cs=11)
section_label(s, '拟合状态', top=1.6)
d2 = [
    ['状态', '训练集', '测试集', '原因', '解决'],
    ['欠拟合', '差', '差', '模型太简单', '增复杂度、加特征'],
    ['过拟合', '极好', '差', '模型太复杂/记噪声', '正则化、加数据、降复杂度'],
    ['泛化好', '好', '好', '平衡', '最终目标'],
]
add_table(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(2.0), d2,
          cw=[Inches(0.9), Inches(0.9), Inches(0.9), Inches(1.5), Inches(1.6)], cs=10)
section_label(s, '正则化 L1 / L2', top=4.2)
d3 = [
    ['类型', '惩罚项 R(w)', '效果', '特点'],
    ['L1 (Lasso)', 'Σ|wⱼ|', '产生稀疏解（很多w=0）', '自动特征选择'],
    ['L2 (Ridge)', 'Σwⱼ²', '整体压小w（不为0）', '平滑稳定，通用首选'],
]
add_table(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(1.5), d3,
          cw=[Inches(1.5), Inches(2.5), Inches(3.5), Inches(4.6)], cs=11)
add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
         'L总 = L数据 + λ·R(w)   |   λ越大惩罚越强', 12, False, GRAY_L)

# ========== Slide 5: 损失函数 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '损失函数：MSE · 交叉熵')
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
         '核心目标：让损失函数最小。损失 = 预测值与真实值的差距。', 13, False, GRAY_L)
d = [
    ['维度', '均方误差 MSE', '交叉熵 Cross-Entropy'],
    ['适用任务', '回归（y连续）', '分类（y离散）'],
    ['公式', 'L = (1/n)Σ(ŷᵢ - yᵢ)²', 'L = -(1/n)Σ[yᵢlogŷᵢ + (1-yᵢ)log(1-ŷᵢ)]'],
    ['直觉', '误差平方平均，大误差重罚', '预测错且自信 → 惩罚极大'],
    ['为什么用', '处处可导，正态分布下最优', 'Sigmoid组合为凸函数，全局最优'],
    ['搭配模型', '线性回归', '逻辑回归、神经网络'],
]
add_table(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(3.6), d,
          cw=[Inches(1.2), Inches(4.5), Inches(6.4)], cs=11)
bullets(s, Inches(0.6), Inches(6.0), Inches(12), Inches(1.2), [
    '分类不用 MSE：Sigmoid+MSE 非凸，多局部最小值；交叉熵+Sigmoid 是凸函数',
    'MSE 用平方：对大误差惩罚更大 + 处处可导便于梯度计算',
])


# ========== Slide 6: 线性回归 vs 逻辑回归 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '线性模型：线性回归 · 逻辑回归')
d = [
    ['维度', '线性回归', '逻辑回归'],
    ['任务', '回归（y连续）', '分类（y离散）'],
    ['模型', 'ŷ = wᵀx + b', 'ŷ = σ(wᵀx + b) = 1/(1+e^(-z))'],
    ['输出', '任意实数', '(0, 1) 概率'],
    ['损失函数', 'MSE', '交叉熵'],
    ['梯度 ∂L/∂w', '(2/n)Σ(ŷᵢ-yᵢ)·xᵢ', '(1/n)Σ(ŷᵢ-yᵢ)·xᵢ'],
    ['决策方式', '直接输出值', 'ŷ > 0.5 判为正类'],
    ['正则化', 'L1 / L2 均可', 'L1 / L2 均可'],
]
add_table(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.5), d,
          cw=[Inches(1.5), Inches(4.8), Inches(5.8)], cs=11)
section_label(s, '梯度下降核心', top=6.3)
add_text(s, Inches(1.8), Inches(6.3), Inches(10), Inches(0.4),
         'w ← w - η·∂L/∂w   |   学习率η太大震荡发散，太小收敛慢', 12)

# ========== Slide 7: KNN vs K-Means ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '距离与聚类：KNN · K-Means')
section_label(s, '距离度量')
d1 = [
    ['名称', '公式', '直觉', '特点'],
    ['欧氏距离', '√Σ(aᵢ-bᵢ)²', '两点直线距离', '最常用，对大差异敏感'],
    ['曼哈顿距离', 'Σ|aᵢ-bᵢ|', '城市格子路', '对异常值不敏感'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.6), d1,
          cw=[Inches(1.2), Inches(1.8), Inches(1.2), Inches(1.3)], cs=10)
add_text(s, Inches(0.6), Inches(3.7), Inches(5.5), Inches(0.4),
         '铁律：距离类算法必须先标准化', 11, True, ACCENT2)
section_label(s, 'KNN vs K-Means', top=1.6)
d2 = [
    ['维度', 'KNN', 'K-Means'],
    ['学习类型', '监督', '无监督'],
    ['K的含义', '最近几个邻居', '分成几个簇'],
    ['需要标签', '需要', '不需要'],
    ['训练', '不训练（存数据）', '迭代找中心'],
    ['预测', '邻居投票/平均', '分配到最近簇'],
    ['K选择', '交叉验证选', '手肘法/轮廓系数'],
]
add_table(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(3.8), d2,
          cw=[Inches(1.2), Inches(2.0), Inches(2.7)], cs=11)
bullets(s, Inches(0.6), Inches(5.8), Inches(12), Inches(1.5), [
    'KNN K太小→过拟合（对噪声敏感），K太大→欠拟合（永远预测多数类）',
    'K-Means 对初始化敏感（K-Means++缓解），假设球形簇，非球形用DBSCAN',
    '共同点：都依赖距离度量，都需要标准化',
])

# ========== Slide 8: SVM ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '支持向量机 SVM')
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
         '核心：找间隔最大的决策边界。只有支持向量（离边界最近的点）决定边界。', 13, False, GRAY_L)
d = [
    ['概念', '说明', '类比'],
    ['支持向量', '离决策边界最近的点，踩在间隔边界上', '撑起帐篷的柱子'],
    ['最大化间隔', '最小化 ||w||，间隔 = 2/||w||', '路越宽越好走'],
    ['软间隔 C', 'C大→严格窄间隔易过拟合；C小→宽松欠拟合', '校规严松'],
    ['核技巧', '映射到高维空间，无需计算高维坐标', '升维找切割面'],
    ['RBF核 γ', 'γ大→边界曲折过拟合；γ小→边界平滑欠拟合', '影响范围大小'],
]
add_table(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(3.6), d,
          cw=[Inches(1.5), Inches(6.5), Inches(4.1)], cs=11)
section_label(s, '核函数对比', top=6.0)
dk = [
    ['核函数', '适用场景'],
    ['线性核', '数据本身线性可分'],
    ['RBF核', '最常用，映射到无穷维'],
    ['多项式核', '需要特征交互'],
]
add_table(s, Inches(0.6), Inches(6.4), Inches(5.0), Inches(0.9), dk,
          cw=[Inches(2.2), Inches(2.8)], cs=10)
add_text(s, Inches(6.5), Inches(6.4), Inches(6), Inches(0.9),
         'vs 逻辑回归：SVM 只看边界关键点（局部）；LR 所有点都影响（全局）',
         11, False, GRAY_L, anchor=MSO_ANCHOR.MIDDLE)

# ========== Slide 9: 决策树 + 随机森林 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '阶段二：树模型与集成', '决策树 · 随机森林')
section_label(s, '决策树 DT')
d1 = [
    ['维度', '说明'],
    ['核心', '每个节点找特征+阈值，使分完后最"纯净"'],
    ['不纯度', '基尼 G=1-Σpₖ²（默认，快） / 信息熵 H=-Σpₖlogpₖ'],
    ['剪枝', '预剪枝(max_depth等) + 后剪枝 → 防过拟合'],
    ['优点', '可解释性极强、不需标准化'],
    ['缺点', '易过拟合、不稳定'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(3.0), d1,
          cw=[Inches(1.2), Inches(4.6)], cs=11)
section_label(s, '随机森林 RF', top=1.6)
d2 = [
    ['维度', '说明'],
    ['本质', 'Bagging + 特征随机 → N棵树投票/平均'],
    ['降方差', 'Var森林 = ρ·σ² + (1-ρ)/N·σ²'],
    ['特征随机', '每次分裂随机选√d个特征 → 降树间相关性ρ'],
    ['优点', '不易过拟合、不需标准化、自带OOB评估'],
    ['缺点', '黑盒不可解释'],
]
add_table(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(3.0), d2,
          cw=[Inches(1.2), Inches(4.6)], cs=11)
section_label(s, '决策树 vs 随机森林', top=5.2)
d3 = [
    ['维度', '决策树', '随机森林'],
    ['结构', '一棵树', 'N棵树投票/平均'],
    ['方差', '高（过拟合）', '低（多树平均降方差）'],
    ['偏差', '低（树可长很深）', '低（和单棵树差不多）'],
    ['可解释性', '强', '弱（黑盒）'],
    ['稳定性', '差', '好'],
]
add_table(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.6), d3,
          cw=[Inches(1.5), Inches(3.5), Inches(7.1)], cs=10)


# ========== Slide 10: Bagging vs Boosting ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '集成两大流派：Bagging · Boosting')
d = [
    ['维度', 'Bagging（随机森林）', 'Boosting（GBDT/XGBoost）'],
    ['训练方式', '并行（树独立）', '串行（后棵纠正前棵）'],
    ['基模型', '深树（低偏差高方差）', '浅树（高偏差低方差）'],
    ['目标', '降方差', '降偏差'],
    ['数据', '每棵树用不同bootstrap抽样', '同一批数据，目标(残差)不同'],
    ['最终输出', '投票/平均', '加权累加'],
    ['过拟合风险', '低', '较高（需早停+正则化）'],
    ['异常值', '不敏感', '敏感（反复放大）'],
]
add_table(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(4.2), d,
          cw=[Inches(1.5), Inches(4.5), Inches(6.1)], cs=11)
bullets(s, Inches(0.6), Inches(6.2), Inches(12), Inches(1.0), [
    'Bagging精髓：深树保低偏差，多树平均降方差 → 不改善偏差，只降低方差',
    'Boosting精髓：串行拟合残差，每棵树走一步梯度下降 → 逐步降偏差',
])

# ========== Slide 11: GBDT + XGBoost + 三大框架 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, 'Boosting 三巨头：GBDT · XGBoost · LightGBM · CatBoost')
section_label(s, 'GBDT 核心')
d1 = [
    ['维度', '说明'],
    ['目标', '拟合残差 = y - ŷ（MSE损失的负梯度）'],
    ['输出', 'ŷ = ŷ₀ + η·树1 + η·树2 + ... + η·树N'],
    ['学习率η', '小→需更多树，但泛化更好'],
    ['本质', '梯度下降 + 决策树，每棵树走一步'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(2.4), d1,
          cw=[Inches(1.2), Inches(4.3)], cs=11)
section_label(s, 'XGBoost vs GBDT', top=1.6)
d2 = [
    ['改进', 'GBDT', 'XGBoost'],
    ['正则化', '无', '损失函数加 L1/L2'],
    ['导数', '一阶', '一阶+二阶（收敛更快）'],
    ['缺失值', '需手动处理', '自动学习走左/右'],
    ['并行', '串行', '特征预排序并行'],
]
add_table(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(2.4), d2,
          cw=[Inches(1.2), Inches(1.8), Inches(2.8)], cs=11)
section_label(s, '三大框架对比', top=4.6)
d3 = [
    ['维度', 'XGBoost', 'LightGBM', 'CatBoost'],
    ['核心优势', '精度高、稳定', '速度快、省内存', '类别特征自动处理'],
    ['类别特征', '需手动编码', '支持但需指定', '全自动'],
    ['生长策略', 'Level-wise', 'Leaf-wise', 'Level-wise（对称树）'],
    ['速度', '中', '最快', '慢（Ordered开销）'],
    ['调参难度', '中', '中', '低（默认就好）'],
]
add_table(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.2), d3,
          cw=[Inches(1.5), Inches(2.8), Inches(3.0), Inches(4.8)], cs=10)

# ========== Slide 12: LightGBM + CatBoost 核心创新 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, 'LightGBM · CatBoost 核心创新')
section_label(s, 'LightGBM 两大创新')
d1 = [
    ['创新', '做法', '效果'],
    ['直方图算法', '连续特征离散化到桶(bin)，分裂只在桶边界尝试', '速度快几倍、内存省4倍，精度损失<1%'],
    ['Leaf-wise生长', '每次只分裂收益最大的叶子（非逐层）', '效率高，但可能过深→需限制num_leaves'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.6), d1,
          cw=[Inches(1.8), Inches(5.5), Inches(4.8)], cs=11)
section_label(s, 'CatBoost 核心：Ordered Target Encoding', top=3.8)
d2 = [
    ['编码方法', '问题'],
    ['One-Hot', '类别多时维度爆炸'],
    ['Label Encoding', '引入虚假大小关系'],
    ['Target Encoding', '数据泄漏：用了全量标签信息'],
    ['Ordered Target Encoding', '每个样本只用"之前"的同类样本算编码 → 无泄漏'],
]
add_table(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.2), d2,
          cw=[Inches(3.0), Inches(9.1)], cs=11)
add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         '选型：类别特征多用CatBoost，大数据用LightGBM，小数据/求稳用XGBoost', 12, True, ACCENT2)

# ========== Slide 13: 模型融合 Stacking/Blending ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '模型融合：Stacking · Blending')
section_label(s, '融合三层次')
d1 = [
    ['层次', '方法', '复杂度', '效果'],
    ['简单融合', '投票（分类）/ 平均（回归）', '最低', '有提升'],
    ['加权融合', '按模型表现分配权重', '低', '稍好'],
    ['Stacking/Blending', '元模型从数据中学习如何组合', '高', '最好'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.8), d1,
          cw=[Inches(2.0), Inches(4.5), Inches(2.0), Inches(3.6)], cs=11)
section_label(s, 'Stacking vs Blending', top=4.0)
d2 = [
    ['维度', 'Stacking', 'Blending'],
    ['元特征生成', 'K折交叉验证（out-of-fold）', 'holdout集（划分一份出来）'],
    ['元模型训练数据', '全部样本', '仅holdout部分'],
    ['数据利用率', '高', '低'],
    ['计算量', '大（K倍训练）', '小（1倍）'],
    ['适用', '数据多、竞赛', '数据少、快速原型'],
]
add_table(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.6), d2,
          cw=[Inches(2.0), Inches(4.5), Inches(5.6)], cs=11)
add_text(s, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
         '关键：基础模型要有多样性（错误模式不同）；元模型越简单越好（逻辑回归即可）', 11, False, GRAY_L)


# ========== Slide 14: EDA + 特征编码 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '阶段三：特征工程与数据处理', 'EDA · 特征编码')
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
         '总纲：模型是锅，特征是米。特征工程决定模型上限，模型只是逼近这个上限。', 13, False, GRAY_L)
section_label(s, 'EDA 五步流程')
d1 = [
    ['步骤', '内容', '类比'],
    ['① 全局概览', 'shape/dtypes/head/describe → 数据规模、类型、基本统计', '称体重、量身高'],
    ['② 目标变量', '分布、偏度 → 决定是否log变换', '先确认病人哪里不舒服'],
    ['③ 单变量分析', '数值型看分布/异常；类别型看占比/失衡', '逐间房看'],
    ['④ 双变量分析', '相关性热力图、散点图 → 找重点特征/共线性', '画人物关系图'],
    ['⑤ 缺失与异常', '缺失类型判断、异常值检测（3σ/IQR）', '专找数据里的坑'],
]
add_table(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(2.8), d1,
          cw=[Inches(1.5), Inches(6.5), Inches(4.1)], cs=11)
section_label(s, '特征编码速查', top=5.1)
d2 = [
    ['编码', '做法', '适用', '坑'],
    ['One-hot', '每类开一列，属于=1否则=0', '无序+类别少(<20)', '类别多→维度爆炸'],
    ['Ordinal', '按序映射1,2,3…', '有序类别（S/M/L）', '无序乱编号=捏造假大小'],
    ['Frequency', '类别→出现次数', '高基数+频率有信息', '同频类别无法区分'],
    ['Target', '类别→该类目标均值', '高基数，榨干信息量', '必须K折防泄漏'],
]
add_table(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.8), d2,
          cw=[Inches(1.5), Inches(3.5), Inches(3.5), Inches(3.6)], cs=10)

# ========== Slide 15: 缺失值 + 异常值 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '缺失值 · 异常值')
section_label(s, '缺失值三种类型')
d1 = [
    ['类型', '含义', '例子', '对策'],
    ['MCAR 完全随机', '纯属手滑', '录入员漏填', '直接补均值/中位数'],
    ['MAR 有条件随机', '缺失与其他特征有关', '豪宅不愿报面积', '用其他特征预测补（KNN插补）'],
    ['MNAR 非随机', '缺失本身携带信息', '没车库→车库面积为空', '不补！缺失即信息，填"None"'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.0), d1,
          cw=[Inches(2.0), Inches(3.0), Inches(3.0), Inches(4.1)], cs=11)
section_label(s, '异常值检测', top=4.2)
d2 = [
    ['方法', '规则', '适用'],
    ['3σ原则', '偏离均值3个标准差之外', '要求数据近正态'],
    ['IQR法', '超出 Q1-1.5×IQR 或 Q3+1.5×IQR', '更稳健，不假设分布'],
]
add_table(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(1.2), d2,
          cw=[Inches(2.0), Inches(5.5), Inches(4.6)], cs=11)
section_label(s, '特征变换对比', top=6.0)
d3 = [
    ['方法', '数据要求', '特点'],
    ['对数变换 log1p', 'x > -1', '最简单，+1兼容x=0'],
    ['Box-Cox', '必须严格正数', '自动寻优λ，λ=0退化为log'],
    ['Yeo-Johnson', '支持零和负数', 'sklearn推荐默认'],
]
add_table(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(1.0), d3,
          cw=[Inches(2.5), Inches(3.5), Inches(6.1)], cs=10)

# ========== Slide 16: 特征构造 + 分箱 + 特征选择 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '特征构造 · 分箱 · 特征选择')
section_label(s, '特征构造套路')
d1 = [
    ['类型', '例子', '直觉'],
    ['加和', 'TotalSF = 地下室+一楼+二楼', '买家看总面积，没人分开算'],
    ['时间差', '房龄 = 销售年份 - 建造年份', '"几岁"比"哪年生"更影响价格'],
    ['比率', '平均每房面积 = 总面积/房间数', '同样100㎡，3房和6房天壤之别'],
    ['乘法交互', '质量分 × 总面积', '好质量让每平米更值钱，协同效应'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.2), d1,
          cw=[Inches(1.5), Inches(4.0), Inches(6.6)], cs=11)
section_label(s, '分箱三种方法', top=4.4)
d2 = [
    ['方法', '做法', '适用'],
    ['等宽 pd.cut', '按值域均分', '分布均匀时'],
    ['等频 pd.qcut', '每箱装同样多样本', '偏态数据首选'],
    ['业务规则', '按领域知识划分', '有领域知识时最香'],
]
add_table(s, Inches(0.6), Inches(4.8), Inches(5.8), Inches(1.8), d2,
          cw=[Inches(1.8), Inches(2.5), Inches(1.5)], cs=11)
section_label(s, '特征选择三流派', top=4.4)
d3 = [
    ['流派', '类比', '做法', '优点', '缺点'],
    ['Filter过滤', '看简历筛', '不看配合，只看单项指标', '极快', '不知岗位匹配度'],
    ['Wrapper包装', '试用期考核', '实际训模型决定去留(RFE)', '最准', '每删一个重训，极慢'],
    ['Embedded嵌入', '末位淘汰', '训练中模型自己打分', '性价比最高', '依赖模型打分机制'],
]
add_table(s, Inches(6.9), Inches(4.8), Inches(5.8), Inches(2.2), d3,
          cw=[Inches(1.2), Inches(1.0), Inches(1.8), Inches(1.0), Inches(0.8)], cs=9)
add_text(s, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
         '实战建议：Filter粗筛 → Embedded精筛；Wrapper只在特征少(<50)时用', 11, False, GRAY_L)


# ========== Slide 17: 类别不平衡 + Pipeline ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '类别不平衡 · Pipeline')
section_label(s, '不平衡处理：SMOTE vs Class Weight')
d1 = [
    ['方法', '做法', '适用场景', '注意事项'],
    ['SMOTE', '在少数类样本间"杂交"造新样本', '少数类极少（几十条）', '只能用在训练集！先SMOTE再划分=泄漏'],
    ['Class Weight', '错杀少数类的惩罚放大N倍', '样本充足、树模型', '一行代码，不造数据'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.5), d1,
          cw=[Inches(1.8), Inches(4.0), Inches(3.0), Inches(3.3)], cs=11)
add_text(s, Inches(0.6), Inches(3.6), Inches(12), Inches(0.4),
         '不平衡数据评估用 PR-AUC / F1，不看 Accuracy（100万里10件危险品，全放行也99.999%）',
         11, True, ACCENT2)
section_label(s, 'Pipeline 四大作用', top=4.2)
d2 = [
    ['作用', '说明', '类比'],
    ['① 防泄漏', 'cross_val_score每折只在训练折fit预处理', '5次模拟考，次次真闭卷'],
    ['② 保一致', '预测自动套训练参数，不会漏步骤/用错模式', '全国门店一个味'],
    ['③ 可搜索', '双下划线穿透子步骤，预处理选择也成超参', '配方本身也能优化'],
    ['④ 易部署', 'joblib.dump(pipe) 预处理+模型一个文件打包', '料理包整包配送'],
]
add_table(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.2), d2,
          cw=[Inches(1.5), Inches(6.0), Inches(4.6)], cs=11)
section_label(s, 'Pipeline vs ColumnTransformer', top=6.9)
add_text(s, Inches(1.8), Inches(6.9), Inches(10.5), Inches(0.4),
         'Pipeline纵向串联（时间轴：先做什么后做什么） | ColumnTransformer横向分流（空间轴：哪列走哪条路）',
         11, False, GRAY_L)

# ========== Slide 18: 模型与特征工程对应关系 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '对症下模型：技术 × 模型 速查')
d = [
    ['技术', '服务谁', '一句话'],
    ['log / Box-Cox / YJ', 'LR、SVM、KNN', '掰正歪分布，树模型免了'],
    ['标准化/归一化', '梯度下降类、距离类', '统一尺度，梯度下降直指圆心'],
    ['特征组合/交互', '所有模型', '通用增益，但要消融验证'],
    ['分箱', '线性模型', '帮直线学会"拐弯"，树模型自己会找切点'],
    ['One-hot编码', '所有模型', '无序类别少的时候用'],
    ['Target编码', '所有模型', '高基数类别，必须K折防泄漏'],
    ['SMOTE', '线性模型、KNN', '少数类极少时用，树模型用Class Weight'],
    ['正则化 L1/L2', 'LR、SVM、NN', '防过拟合，L1还能做特征选择'],
]
add_table(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(4.5), d,
          cw=[Inches(2.5), Inches(3.5), Inches(6.1)], cs=11)
bullets(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.8), [
    '树模型不需要特征标准化/变换 → 按排序找切分点，单调变换排序不变',
    '但对目标y做log变换对树模型依然有用 → MSE不再被天价豪宅主导',
])

# ========== Slide 19: 总结 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '学习总结')

summary = [
    ('阶段一 基础原理', [
        'ML ⊃ NN ⊃ DL；监督有标签，无监督找结构',
        '参数模型学，超参数人定；过拟合降方差，欠拟合增复杂度',
        '线性回归用MSE，逻辑回归用交叉熵；梯度下降沿负梯度逐步逼近',
        'KNN看邻居投票，K-Means迭代找中心；都依赖距离，必须标准化',
        'SVM找最大间隔，支持向量决定边界；核技巧升维，RBF最常用',
    ]),
    ('阶段二 树模型与集成', [
        '决策树可解释强但易过拟合；随机森林Bagging降方差',
        'Boosting串行降偏差；GBDT拟合残差，XGBoost全面优化',
        'LightGBM直方图+Leaf-wise更快更省内存',
        'CatBoost Ordered Target Encoding自动处理类别特征无泄漏',
        'Stacking用元模型学组合方式，K折生成元特征防泄漏',
    ]),
    ('阶段三 特征工程', [
        'EDA五步走：全局→目标→单变量→双变量→缺失异常',
        '编码四法：One-hot/Ordinal/Frequency/Target（Target必K折）',
        '变换：log/Box-Cox/YJ 掰正分布；树模型不需要但y变换仍有用',
        '特征选择：Filter粗筛→Embedded精筛→Wrapper少用',
        'Pipeline+ColumnTransformer：防泄漏、保一致、易部署',
    ]),
]

y_start = 1.7
for i, (title, items) in enumerate(summary):
    col_x = 0.6 + i * 4.2
    add_text(s, Inches(col_x), Inches(y_start), Inches(3.8), Inches(0.4),
             title, 15, True, ACCENT)
    bullets(s, Inches(col_x), Inches(y_start + 0.5), Inches(3.8), Inches(4.5),
            items, size=10, color=WHITE)

# 底部金句
add_text(s, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.5),
         '特征工程决定模型上限，模型只是逼近这个上限',
         14, True, ACCENT2, align=PP_ALIGN.CENTER)

# ========== Slide 20: 结束页 ==========
s = prs.slides.add_slide(BL); set_bg(s)

for y in [3.4, 3.45, 3.5]:
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(y), Inches(2.3), Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()

add_text(s, Inches(0), Inches(2.5), Inches(13.333), Inches(1.0),
         'THANKS', 72, True, WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0), Inches(3.8), Inches(13.333), Inches(0.6),
         '感谢聆听', 24, False, GRAY_L, align=PP_ALIGN.CENTER)

add_text(s, Inches(0), Inches(6.5), Inches(13.333), Inches(0.4),
         '暑期培训 · 机器学习方向 · 2026', 12, False, GRAY_M, align=PP_ALIGN.CENTER)


# ========== 保存 ==========
output_path = '/home/xavier/暑期培训/02 机器学习/机器学习学习汇报.pptx'
prs.save(output_path)
print(f'PPT已生成: {output_path}')
print(f'共 {len(prs.slides)} 页')
