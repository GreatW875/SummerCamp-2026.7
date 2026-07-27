#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""机器学习学习汇报PPT生成脚本(下) | Ubuntu 22.04 + conda(sc) + python-pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY_D = RGBColor(0x33, 0x33, 0x33)
GRAY_M = RGBColor(0x66, 0x66, 0x66)
GRAY_L = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0, 0xB4, 0xD8)
ACCENT2 = RGBColor(0xFF, 0x6B, 0x6B)
TH_BG = RGBColor(0x1A, 0x1A, 0x1A)
TR_BG = RGBColor(0x0D, 0x0D, 0x0D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BL = prs.slide_layouts[6]


def set_bg(s, c=BLACK):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()


def _set_font(run, size, bold, color, font='微软雅黑'):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)


def add_text(s, l, t, w, h, text, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    _set_font(run, size, bold, color)
    return tb


def title_bar(s, title, sub=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.08), Inches(0.45))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    add_text(s, Inches(0.85), Inches(0.45), Inches(11), Inches(0.6), title, 28, True)
    if sub:
        add_text(s, Inches(0.85), Inches(0.95), Inches(11), Inches(0.35), sub, 13, False, GRAY_L)
    sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.1), Emu(9525))
    sep.fill.solid(); sep.fill.fore_color.rgb = GRAY_D; sep.line.fill.background()


def section_label(s, text, top=1.6):
    add_text(s, Inches(0.6), Inches(top), Inches(5), Inches(0.35), text, 12, True, ACCENT)


def _set_cell_border(cell, color='333333'):
    tcPr = cell._tc.get_or_add_tcPr()
    for bn in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = tcPr.find(qn(bn))
        if ln is None:
            ln = etree.SubElement(tcPr, qn(bn))
        ln.set('w', '6350')
        sf = ln.find(qn('a:solidFill'))
        if sf is None:
            sf = etree.SubElement(ln, qn('a:solidFill'))
        sc = sf.find(qn('a:srgbClr'))
        if sc is None:
            sc = etree.SubElement(sf, qn('a:srgbClr'))
        sc.set('val', color)


def add_table(s, l, t, w, h, data, cw=None, hs=12, cs=11, fcb=True):
    rows, cols = len(data), len(data[0])
    ts = s.shapes.add_table(rows, cols, l, t, w, h)
    tbl = ts.table
    if cw:
        for i, wi in enumerate(cw):
            tbl.columns[i].width = wi
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c); cell.text = ''
            tf = cell.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_right = Emu(45720)
            tf.margin_top = tf.margin_bottom = Emu(22860)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = str(data[r][c])
            is_h = (r == 0)
            is_fcb = fcb and c == 0 and r > 0
            sz = hs if is_h else cs
            bold = is_h or is_fcb
            clr = ACCENT if is_h else (ACCENT2 if is_fcb else WHITE)
            _set_font(run, sz, bold, clr)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TH_BG if is_h else (TR_BG if r % 2 == 1 else BLACK)
            _set_cell_border(cell)
    return tbl


def bullets(s, l, t, w, h, items, size=12, color=GRAY_L):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(5)
        run = p.add_run(); run.text = f'▸  {item}'
        _set_font(run, size, False, color)
    return tb


# ========== Slide 1: 封面 ==========
s = prs.slides.add_slide(BL); set_bg(s)
for y in [2.2, 2.25, 2.3]:
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(y), Inches(0.6), Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
add_text(s, Inches(1.5), Inches(2.6), Inches(10), Inches(1.2), '机器学习学习汇报（下）', 48, True)
add_text(s, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
         '模型评估诊断调优 · 项目实战与部署', 22, False, GRAY_L)
add_text(s, Inches(1.5), Inches(6.2), Inches(6), Inches(0.4), '暑期培训 · 机器学习方向', 14, False, GRAY_M)
add_text(s, Inches(1.5), Inches(6.6), Inches(6), Inches(0.4), '2026', 12, False, GRAY_D)
vl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.5), Inches(2.5), Emu(9525), Inches(3))
vl.fill.solid(); vl.fill.fore_color.rgb = GRAY_D; vl.line.fill.background()

# ========== Slide 2: 目录 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '目录 / CONTENTS')
toc = [
    ('01', '评估指标全家桶', '混淆矩阵 · P/R/F1 · ROC-AUC · PR-AUC · KS'),
    ('02', '交叉验证与数据泄漏', 'KFold · StratifiedKFold · TimeSeriesSplit · 泄漏四大场景'),
    ('03', '学习曲线与验证曲线', '高偏差/高方差诊断 · 学习曲线 · 验证曲线'),
    ('04', '调优与解释', 'Optuna贝叶斯调参 · 概率校准 · SHAP模型解释'),
    ('05', '项目实战与部署', '完整流程 · Joblib · FastAPI · Docker · 文档'),
]
for i, (n, t, d) in enumerate(toc):
    y = 1.8 + i * 1.1
    add_text(s, Inches(1.0), Inches(y), Inches(1.0), Inches(0.6), n, 36, True, ACCENT)
    add_text(s, Inches(2.3), Inches(y + 0.02), Inches(8), Inches(0.4), t, 20, True)
    add_text(s, Inches(2.3), Inches(y + 0.45), Inches(9), Inches(0.35), d, 12, False, GRAY_L)
    if i < len(toc) - 1:
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(y + 0.95), Inches(11), Emu(6350))
        sp.fill.solid(); sp.fill.fore_color.rgb = GRAY_D; sp.line.fill.background()

# ========== Slide 3: 混淆矩阵 + 核心指标 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '阶段四：模型评估、诊断与调优', '混淆矩阵 · 核心指标')
section_label(s, '混淆矩阵（癌症筛查为例）')
d1 = [
    ['', '预测：有病', '预测：健康'],
    ['实际：有病', 'TP 真阳性 ✅ 抓对了', 'FN 假阴性 ❌ 漏诊'],
    ['实际：健康', 'FP 假阳性 ❌ 误诊', 'TN 真阴性 ✅ 放对了'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(1.8), d1,
          cw=[Inches(1.8), Inches(2.2), Inches(1.8)], cs=11)
section_label(s, '四大核心指标', top=1.6)
d2 = [
    ['指标', '公式', '回答的问题', '类比'],
    ['Accuracy', '(TP+TN)/全部', '整体猜对多少？', '考试总分'],
    ['Precision', 'TP/(TP+FP)', '我说有病的里面多少是真的？', '不乱冤枉人'],
    ['Recall', 'TP/(TP+FN)', '真患者里我揪出了多少？', '不漏网'],
    ['F1', 'P与R的调和平均', 'P和R的综合分', '不偏科的综合成绩'],
]
add_table(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(2.4), d2,
          cw=[Inches(1.2), Inches(1.8), Inches(2.0), Inches(0.8)], cs=10)
add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.4),
         '关键直觉：Precision 和 Recall 是跷跷板。调松→R涨P跌；调严→反过来。', 12, True, ACCENT2)
section_label(s, '为什么F1用调和平均', top=4.5)
add_text(s, Inches(0.6), Inches(4.9), Inches(12), Inches(0.4),
         '调和平均被短板狠狠拉低——P=100%、R=2%时普通平均还有51%，F1只有约4%，专治"瘸腿"模型。',
         12, False, WHITE)
section_label(s, '指标速查表', top=5.5)
d3 = [
    ['指标', '本质', '依赖阈值？', '不平衡适用', '典型场景'],
    ['Accuracy', '整体猜对率', '是', '❌ 严重误导', '均衡数据粗看'],
    ['Precision', '报警可信度', '是', '✅', '误报代价高（垃圾邮件）'],
    ['Recall', '抓全能力', '是', '✅', '漏报代价高（癌症、欺诈）'],
    ['F1', 'P/R综合', '是', '✅', '不平衡但要单值总结'],
    ['ROC-AUC', '排序区分能力', '否', '⚠️ 虚高', '均衡/中等不平衡'],
    ['PR-AUC', '少数派上的排序', '否', '✅✅ 首选', '极端不平衡'],
    ['KS', '两条累计曲线最大落差', '否', '✅', '金融风控评分卡'],
]
add_table(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.5), d3,
          cw=[Inches(1.2), Inches(2.0), Inches(1.5), Inches(1.5), Inches(5.9)], cs=9)


# ========== Slide 4: ROC-AUC vs PR-AUC + KS ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '排序能力指标：ROC-AUC · PR-AUC · KS')
section_label(s, 'ROC-AUC')
d1 = [
    ['维度', '说明'],
    ['ROC', '阈值从1扫到0，画(真正例率, 假正例率)曲线'],
    ['AUC', 'ROC曲线下面积；AUC=0.5瞎猜，越接近1越好'],
    ['直觉', '随机抽一正一负，模型把正的排更高的概率'],
    ['类比', '考试排名能力'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.2), d1,
          cw=[Inches(1.2), Inches(4.6)], cs=11)
section_label(s, 'PR-AUC：不平衡照妖镜', top=1.6)
d2 = [
    ['维度', '说明'],
    ['问题', 'ROC横轴是假正例率，健康人基数大→误诊100个比例仍很小→曲线虚高'],
    ['PR曲线', '直接画Precision vs Recall，只盯着少数派看'],
    ['选择', '均衡→ROC-AUC；正样本极少（欺诈/癌症/违约）→PR-AUC'],
]
add_table(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(2.2), d2,
          cw=[Inches(1.2), Inches(4.6)], cs=11)
section_label(s, 'KS值：风控专用', top=4.4)
d3 = [
    ['维度', '说明'],
    ['做法', '按模型分数排序，累计"抓到的患者比例"和"误伤的健康人比例"，两条曲线最大距离'],
    ['类比', '筛金子——按含金量过筛，最好的时候能甩开沙子多远'],
    ['经验值', 'KS>0.3可用，>0.4不错；>0.5要警惕数据泄漏'],
]
add_table(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.8), d3,
          cw=[Inches(1.2), Inches(10.9)], cs=11)
add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
         '一句话：阈值变混淆矩阵变 → 需要不依赖阈值的排序能力指标', 12, True, ACCENT2)

# ========== Slide 5: 交叉验证 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '交叉验证的正确姿势')
section_label(s, '为什么需要交叉验证')
add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.4),
         '留出法局限：一次估分运气成分大。K折：轮流考K次取平均 → 分数更稳、数据不浪费、能看波动', 12, False, WHITE)
section_label(s, '三种交叉验证对比', top=2.6)
d = [
    ['方法', '切法核心', '类比', '什么时候用'],
    ['KFold', '随机切K份轮流考', '5份卷子轮着模考', '回归问题默认'],
    ['StratifiedKFold', '每折类别比例配平', '每组都是班级的缩影', '分类问题默认，尤其不平衡'],
    ['TimeSeriesSplit', '只用过去预测未来', '模拟炒股不能穿越', '任何带时间顺序的数据'],
]
add_table(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(2.0), d,
          cw=[Inches(2.2), Inches(3.5), Inches(3.0), Inches(3.4)], cs=11)
section_label(s, 'StratifiedKFold 解决什么问题', top=5.2)
add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
         '癌症筛查患者只占1%，随机切可能某折没患者→白考。分层分保证每折都是1%，和原始分布一致。',
         12, False, WHITE)
section_label(s, 'TimeSeriesSplit 解决什么问题', top=6.2)
add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         '股价/销量/气温有时间顺序，随机切=用明天预测昨天→穿越作弊。永远只用过去预测未来。',
         12, False, WHITE)

# ========== Slide 6: 数据泄漏 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '数据泄漏专题：本阶段第一大坑')
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
         '定义：训练中模型偷看到了本不该知道的信息 → 离线分数虚高，上线立刻翻车。不报错、不报警，给你虚假自信。',
         12, False, GRAY_L)
section_label(s, '四大经典泄漏场景')
d = [
    ['#', '泄漏场景', '泄漏的信息', '对策'],
    ['1', '全数据先预处理再划分', '测试集的分布（均值/方差等）', '预处理放进Pipeline，一切在CV内部发生'],
    ['2', 'Target Encoding不做K折', '测试集的y（答案）', '编码嵌套进交叉验证（out-of-fold）'],
    ['3', '时序数据随机划分', '未来的信息', 'TimeSeriesSplit，只用过去预测未来'],
    ['4', '特征含"答案的影子"', '结果本身（事后产物）', '问"预测时刻拿得到吗"，拿不到就删'],
]
add_table(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(2.8), d,
          cw=[Inches(0.5), Inches(3.5), Inches(3.5), Inches(4.6)], cs=11)
section_label(s, '泄漏自查', top=5.2)
bullets(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.5), [
    '一句话总结：训练过程能接触到的任何信息，必须在真实预测时刻也能拿到',
    '万能自查问题："我做这一步时，用到的数据里有没有混进考试时还不知道的东西？"',
    '危险信号：分数好得反常（AUC 0.99、KS > 0.5），先怀疑泄漏，再庆祝',
    '标准解法：所有预处理塞进 Pipeline，cross_val_score(pipe, X, y, cv=5) 自动防泄漏',
], size=11, color=WHITE)


# ========== Slide 7: 高偏差 vs 高方差 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '学习曲线：先诊断，再开药')
section_label(s, '高偏差 vs 高方差')
d = [
    ['', '高偏差（欠拟合）', '高方差（过拟合）'],
    ['本质', '模型太"笨"，规律没学到', '模型太"轴"，连噪音都背下来了'],
    ['训练集表现', '❌ 差', '✅ 极好'],
    ['验证集表现', '❌ 差', '❌ 差'],
    ['健身类比', '每天只散步，肌肉没刺激到', '背下动作细节，换健身房就不会练了'],
    ['考试类比', '教材都没读懂，模考高考都砸', '背下真题答案，新题全不会'],
    ['药方方向', '加特征、换复杂模型、减正则化', '加数据、降复杂度、加正则化'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(3.5), d,
          cw=[Inches(1.5), Inches(4.8), Inches(5.8)], cs=11)
add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
         '一句话：高偏差是"学不够"，高方差是"学死了"。两者药方完全相反 → 必须先诊断再动手。',
         13, True, ACCENT2)
section_label(s, '判断口诀', top=6.3)
add_text(s, Inches(1.8), Inches(6.3), Inches(10), Inches(0.4),
         '两线都低 → 治笨（加容量）；两线分叉 → 治轴（加数据/降复杂度）', 12, False, WHITE)

# ========== Slide 8: 学习曲线 + 验证曲线 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '学习曲线 · 验证曲线')
section_label(s, '学习曲线 Learning Curve')
d1 = [
    ['维度', '说明'],
    ['横轴', '训练数据量'],
    ['纵轴', '分数（训练分 + 验证分两条线）'],
    ['回答问题', '"再加数据/加特征，还有救吗？"'],
    ['高偏差型', '两线收敛在一起但都低 → 加数据没用，加容量'],
    ['高方差型', '训练分很高、验证分明显低，中间有鸿沟 → 加数据有效'],
    ['工具', 'learning_curve()'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(3.8), d1,
          cw=[Inches(1.5), Inches(4.3)], cs=11)
section_label(s, '验证曲线 Validation Curve', top=1.6)
d2 = [
    ['维度', '说明'],
    ['横轴', '某个超参数的取值（如max_depth、C）'],
    ['纵轴', '分数（训练分 + 验证分两条线）'],
    ['回答问题', '"这个旋钮拧到哪个位置最合适？"'],
    ['左端', '参数太保守 → 两线都低 → 欠拟合区'],
    ['右端', '参数太激进 → 训练分冲天、验证分掉头 → 过拟合区'],
    ['甜蜜点', '验证分峰值位置 → 该选的超参值'],
    ['工具', 'validation_curve()'],
]
add_table(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.2), d2,
          cw=[Inches(1.5), Inches(4.3)], cs=11)
section_label(s, '两张曲线分工', top=6.4)
add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
         '先看学习曲线定大方向 → 再看验证曲线微调旋钮；诊断错了，努力全废。', 12, True, ACCENT2)

# ========== Slide 9: 调参方法对比 + Optuna ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '调优：Optuna 贝叶斯调参')
section_label(s, '三种调参方法对比')
d1 = [
    ['方法', '做法', '问题'],
    ['GridSearch 网格搜索', '每个参数候选值排列组合，全部试一遍', '参数一多组合爆炸：5参数各10值=10万次训练'],
    ['RandomSearch 随机搜索', '随机挑组合试', '不爆炸了，但完全不长记性——烂区域反复踩'],
    ['贝叶斯优化（Optuna）', '每试一组更新猜想地图，下一组挑最可能有惊喜的试', '需要合理设置搜索空间'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.8), d1,
          cw=[Inches(2.5), Inches(5.5), Inches(4.1)], cs=11)
section_label(s, 'Optuna 两大杀手锏', top=4.0)
d2 = [
    ['杀手锏', '说明', '类比'],
    ['TPE采样', '智能多试"高分区域"，少碰"烂区域"', '调收音机：听到杂音强就知道电台在附近，细细微调'],
    ['Pruning剪枝', '某次试验中途明显很差就提前掐死，不浪费算力', '厨师尝一口汤不行直接倒掉，不用做完再倒'],
]
add_table(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(1.4), d2,
          cw=[Inches(2.0), Inches(6.5), Inches(3.6)], cs=11)
section_label(s, '实战工作流', top=6.0)
bullets(s, Inches(0.6), Inches(6.4), Inches(12), Inches(1.0), [
    '① 画学习曲线诊断方向（高偏差/高方差）',
    '② 大方向对了，画验证曲线逐个超参找甜蜜点',
    '③ 最后用 Optuna 系统搜索（≥50 trials）',
    '调参是最后一步不是第一步 —— 方向错了，调参白忙',
], size=11, color=WHITE)


# ========== Slide 10: 概率校准 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '概率校准：修好体重秤')
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
         '模型擅长排序（谁比谁风险高），不擅长报数（准确的百分比）。像选秀评委：排名清楚，百分制乱打。',
         12, False, GRAY_L)
section_label(s, '什么时候需要校准')
d1 = [
    ['场景', '需要校准？'],
    ['只按风险排名圈人、阈值一刀切、比赛只比AUC', '❌ 只用排序，不用校准'],
    ['银行按违约概率定利率、医院按患病概率定手术、保险定保费', '✅ 数字直接换算成钱/人命，必须校准'],
]
add_table(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.4), d1,
          cw=[Inches(9.0), Inches(3.1)], cs=11)
section_label(s, '两种校准方法')
d2 = [
    ['方法', '原理', '类比', '适用'],
    ['Platt Scaling', '假设偏差是S形，用逻辑回归拟合映射', '体重秤固定多报5斤：装个"减5斤"弹簧', '数据量少时稳（仅2参数，不易过拟合）'],
    ['Isotonic Regression', '不假设形状，只要求"分高真实概率不能更低"，拟合阶梯形曲线', '老师傅逐刻度手工校秤，一格一格掰准', '数据量大时更准（数据少容易学歪）'],
]
add_table(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(1.6), d2,
          cw=[Inches(2.2), Inches(4.5), Inches(3.0), Inches(2.4)], cs=11)
section_label(s, '校准到什么程度', top=5.5)
bullets(s, Inches(0.6), Inches(5.9), Inches(12), Inches(1.5), [
    '看校准曲线：预测概率分桶，画"预测vs真实"，大致贴住对角线即可',
    '看ECE（期望校准误差）：各桶预测与真实差值加权平均，0.01~0.05内够用',
    '最终标准是业务：偏差小到不影响决策质量就停',
    '校准不会破坏排序能力——单调映射，不改排名',
], size=11, color=WHITE)

# ========== Slide 11: SHAP 模型解释 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, 'SHAP 模型解释：侦探分功劳')
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
         'SHAP = SHapley Additive exPlanations，回答"这次预测里每个特征各出了多少力"。理论来源：博弈论沙普利值。',
         12, False, GRAY_L)
section_label(s, '三大优点')
d1 = [
    ['优点', '说明'],
    ['可解释到单次预测（局部）', '"为什么这个客户被判违约"有直接答案'],
    ['可加总', '所有特征贡献 + 基准值 = 最终预测，账目严丝合缝'],
    ['可汇总看全局', '几千样本贡献叠一起画summary图，看整体重要性和方向'],
]
add_table(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.8), d1,
          cw=[Inches(3.0), Inches(9.1)], cs=11)
section_label(s, '两种常用图')
d2 = [
    ['图', '看什么', '回答的问题'],
    ['summary_plot（蜂群图）', '特征重要性总排名 + 正负方向', '模型整体靠哪些特征做判断？'],
    ['waterfall_plot（瀑布图）', '单个样本的贡献分解', '为什么这一个客户被判违约？'],
]
add_table(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(1.4), d2,
          cw=[Inches(3.0), Inches(4.5), Inches(4.6)], cs=11)
section_label(s, '直觉类比', top=5.7)
add_text(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
         '4人组队送外卖赚100元小费怎么分？试所有组合顺序，算每人"加入那一刻"平均多带来多少收入。不偏袒、可复算、加起来正好100元。',
         12, False, WHITE)
section_label(s, '附加价值', top=6.7)
add_text(s, Inches(1.8), Inches(6.7), Inches(10), Inches(0.4),
         '兼作数据泄漏探测器：ID类特征排第一 = 警报（特征里混进了答案）',
         11, True, ACCENT2)

# ========== Slide 12: 调优解释一条线 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '阶段四总结：诊断 → 调优 → 校准 → 解释')
d = [
    ['步骤', '工具/方法', '回答的问题'],
    ['① 诊断方向', '学习曲线', '该加数据还是该换脑子？（高偏差/高方差）'],
    ['② 系统调参', 'Optuna 贝叶斯优化', '最优参数组合是什么？（≥50 trials）'],
    ['③ 概率校准', 'Platt / Isotonic', '输出概率可信吗？（拿数字做决策才需要）'],
    ['④ 模型解释', 'SHAP', '为什么这么预测？每个特征贡献多少？'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.4), d,
          cw=[Inches(2.0), Inches(3.5), Inches(6.6)], cs=12)
section_label(s, '完整工作流', top=4.6)
bullets(s, Inches(0.6), Inches(5.0), Inches(12), Inches(2.0), [
    '基线模型 → 学习曲线诊断（高偏差/高方差）→ 调整方向（加特征/加数据/调复杂度）',
    '→ 验证曲线逐个超参找甜蜜点 → Optuna 系统搜索最优组合',
    '→ 概率校准（需要的话） → SHAP 检查合理性 + 错误分析归因',
    '→ 回流特征/数据迭代 → 一个能交代、能上线、能答辩的模型',
], size=12, color=WHITE)


# ========== Slide 13: 阶段五 完整项目流程 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '阶段五：项目实战与部署', '完整项目流程')
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
         '需求 → 数据 → 特征 → 模型 → 评估  |  类比开餐馆：定菜品→买食材→洗菜切配→下锅烹饪→试吃验收',
         12, False, GRAY_L)
d = [
    ['流程', '干什么', '对应知识点', '头号大坑'],
    ['需求', '把"人话"翻译成"机器话"：定y、定分类/回归、定指标', '问题定义、评估指标选取', '方向错，后面全白干'],
    ['数据', '收集 → EDA体检 → 划分', 'EDA、缺失值插补、异常值、泄漏防护', '泄漏——预处理必须进Pipeline'],
    ['特征', '把原料处理成"能下锅"的状态', '编码、变换、构造、特征选择、Pipeline', '只换模型不造特征（特征决定上限）'],
    ['模型', '先基线打底，再升级', 'LR基线、RF、XGBoost、LightGBM、Stacking', '盲目调参（调参是最后一步）'],
    ['评估', '打分 + 诊断 + 归因', '混淆矩阵、AUC、交叉验证、学习曲线、错误分析', '不平衡数据只看Accuracy'],
]
add_table(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(3.5), d,
          cw=[Inches(1.0), Inches(3.5), Inches(4.0), Inches(3.6)], cs=11)
section_label(s, '核心认知', top=5.8)
add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
         '不是单向流水线，而是循环——评估结论会把你打回前面某一步，迭代几圈才收敛。上线后监控漂移，形成完整闭环。',
         12, True, ACCENT2)

# ========== Slide 14: 模型优化与解释 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '模型优化与解释：调参 · SHAP · 错误分析')
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
         '三件事回答三个问题：模型还能不能更好（调参）、为什么这么预测（SHAP）、在哪里会犯错（错误分析）',
         12, False, GRAY_L)
d = [
    ['手段', '核心思想', '类比', '关键要点'],
    ['Optuna调参', '根据上次结果，聪明地选下一组参数', '调洗澡水温，不是瞎拧每个刻度', '诊断先行：学习曲线判方向；≥50 trials、固定random_state'],
    ['SHAP解释', '算每个特征对本次预测的贡献（可正可负）', '团队分奖金：看你"在与不在"差多少', '全局看summary蜂群图、局部看waterfall图；兼作泄漏探测器'],
    ['错误分析', '挑出错样本人工归因', '学霸的错题本', '三方向：特征不够→回特征工程；噪声/标注错→洗数据；边界模糊→接受天花板'],
]
add_table(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(2.6), d,
          cw=[Inches(1.5), Inches(4.0), Inches(3.0), Inches(3.6)], cs=11)
section_label(s, '工作流', top=5.0)
add_text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
         '基线 → 学习曲线诊断 → Optuna调参 → SHAP检查合理性 → 错误分析归因 → 回流特征/数据 → 再迭代',
         13, True, WHITE)
section_label(s, '产出物', top=6.2)
add_text(s, Inches(1.8), Inches(6.2), Inches(10), Inches(0.4),
         '一页模型报告：指标 + 学习曲线 + SHAP图 + 错误样本分析', 12, False, GRAY_L)

# ========== Slide 15: Joblib + FastAPI ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '持久化与服务化：Joblib · FastAPI')
add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
         'Joblib解决"模型怎么存"，FastAPI解决"别人怎么用"。训练=培养厨师，Joblib=冷冻手艺，FastAPI=开点餐窗口。',
         12, False, GRAY_L)
section_label(s, 'Joblib 持久化')
bullets(s, Inches(0.6), Inches(2.1), Inches(5.8), Inches(2.0), [
    '核心两行：joblib.dump(model, "model.pkl") 存 / joblib.load 取',
    '存的是整个Pipeline（预处理+模型打包），保证线上线下一致',
    '选Joblib而非Pickle：对大NumPy数组存取更快',
], size=11, color=WHITE)
section_label(s, 'FastAPI 服务化三零件', top=1.6)
d = [
    ['零件', '干什么', '类比'],
    ['@app.post("/predict")', '定义接口地址', '点餐窗口'],
    ['Pydantic模型 Input', '校验JSON字段与类型，不合规返422', '门口安检，脏数据进不了厨房'],
    ['返回 {"prediction": ...}', '输出预测结果', '出菜'],
]
add_table(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(2.0), d,
          cw=[Inches(2.5), Inches(2.5), Inches(0.8)], cs=11)
section_label(s, '完整链路', top=4.2)
add_text(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.4),
         '训练 → dump存Pipeline → API启动时load一次（常驻内存）→ 收JSON → Pydantic安检 → 转DataFrame → predict → 返回JSON',
         12, False, WHITE)
section_label(s, '性能要点', top=5.3)
add_text(s, Inches(1.8), Inches(5.3), Inches(10), Inches(0.4),
         '模型在服务启动时load一次，不是每来一单才解冻——开门前备好菜，营业时只管炒。', 11, True, ACCENT2)
section_label(s, '附赠亮点', top=6.0)
add_text(s, Inches(1.8), Inches(6.0), Inches(10), Inches(0.4),
         '访问 /docs 自动生成交互式文档页面，可现场试点', 11, False, GRAY_L)
section_label(s, '联调测试', top=6.7)
add_text(s, Inches(1.8), Inches(6.7), Inches(10), Inches(0.4),
         'fastapi dev main.py 起服务 → curl / Postman发JSON → 核对返回', 11, False, GRAY_L)


# ========== Slide 16: Docker + 文档 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, 'Docker容器化 · 项目文档')
section_label(s, 'Docker：把整间厨房打包寄走')
d1 = [
    ['概念', '说明'],
    ['解决痛点', '"在我电脑上明明能跑啊！"——把代码+环境整体打包成集装箱，任何机器即用'],
    ['三件套', 'Dockerfile（图纸）→ build → 镜像Image（样板房）→ run → 容器Container（营业中）'],
    ['两条命令', 'docker build -t 名字 . 造镜像 → docker run -p 8000:8000 名字 起容器'],
    ['端口映射', '大堂8000号窗口对应后厨8000号灶台'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.4), d1,
          cw=[Inches(2.0), Inches(10.1)], cs=11)
section_label(s, 'Dockerfile 六行骨架', top=4.6)
code_lines = [
    'FROM python:3.10-slim        # 毛坯房：带Python的精简系统',
    'WORKDIR /app                 # 进屋定工作目录',
    'COPY requirements.txt .      # 先搬清单（利用构建缓存）',
    'RUN pip install -r requirements.txt   # 按清单装修',
    'COPY . .                     # 再搬家具（全部代码）',
    'CMD ["uvicorn", "main:app", "--host", "0.0.0.0", "--port", "8000"]  # 点火开业',
]
tb = s.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = Emu(91440); tf.margin_right = Emu(91440)
tf.margin_top = Emu(45720); tf.margin_bottom = Emu(45720)
for i, line in enumerate(code_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(2)
    run = p.add_run(); run.text = line
    _set_font(run, 10, False, ACCENT, font='Consolas')
# 代码框背景
code_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.0))
code_bg.fill.solid(); code_bg.fill.fore_color.rgb = TH_BG
code_bg.line.color.rgb = GRAY_D
code_bg.shadow.inherit = False
tb.element.getparent().remove(tb.element)
code_bg._element.addnext(tb.element)

section_label(s, 'README五要素', top=7.2)
add_text(s, Inches(1.8), Inches(7.2), Inches(10), Inches(0.3),
         '项目说明 → 快速开始 → API示例 → 项目结构 → 模型性能（指标+SHAP图）',
         11, False, GRAY_L)

# ========== Slide 17: 阶段五总结 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '阶段五总结：从训练到交付')
d = [
    ['模块', '核心内容', '产出/能力'],
    ['流程串讲', '需求→数据→特征→模型→评估，循环迭代', '完整项目地图，知道每步干什么、坑在哪'],
    ['优化解释', 'Optuna调参 + SHAP解释 + 错误分析', '模型打磨到最优，决策依据说得清'],
    ['持久化服务化', 'Joblib存Pipeline + FastAPI接口', '模型能被别人调用，不是躺在笔记本里'],
    ['容器化文档', 'Docker打包 + README五要素', '别人10分钟能复现，能写进简历'],
]
add_table(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.8), d,
          cw=[Inches(2.0), Inches(5.5), Inches(4.6)], cs=12)
section_label(s, '完整闭环', top=5.0)
add_text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
         '流程串讲（地图）→ 优化解释（打磨）→ 持久化+API（上线）→ Docker+文档（交付）',
         14, True, WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.2), Inches(13.333), Inches(0.5),
         '= 一个能写进简历、能向面试官讲透的完整作品',
         16, True, ACCENT2, align=PP_ALIGN.CENTER)

# ========== Slide 18: 全阶段总结 ==========
s = prs.slides.add_slide(BL); set_bg(s)
title_bar(s, '机器学习全阶段总结')

stages = [
    ('阶段一 基础原理', [
        'ML ⊃ NN ⊃ DL',
        '监督有标签 / 无监督找结构',
        '线性回归MSE / 逻辑回归交叉熵',
        'KNN看邻居 / K-Means找中心',
        'SVM最大间隔 / 核技巧升维',
    ]),
    ('阶段二 树模型集成', [
        '决策树可解释但易过拟合',
        'Bagging降方差（RF）',
        'Boosting降偏差（GBDT/XGB）',
        'LightGBM快 / CatBoost自动类别',
        'Stacking元模型学组合',
    ]),
    ('阶段三 特征工程', [
        'EDA五步走验房',
        '编码四法（Target必K折）',
        '变换掰正分布',
        '构造交互1+1>2',
        'Pipeline防泄漏',
    ]),
    ('阶段四 评估调优', [
        '混淆矩阵 + P/R/F1',
        'ROC-AUC / PR-AUC / KS',
        'K折/分层/时序交叉验证',
        '学习曲线诊断偏差方差',
        'Optuna调参 + SHAP解释',
    ]),
    ('阶段五 实战部署', [
        '需求→数据→特征→模型→评估',
        '基线→诊断→调参→解释',
        'Joblib持久化Pipeline',
        'FastAPI服务化接口',
        'Docker打包 + README文档',
    ]),
]

for i, (title, items) in enumerate(stages):
    col_x = 0.4 + i * 2.55
    add_text(s, Inches(col_x), Inches(1.7), Inches(2.4), Inches(0.4),
             title, 12, True, ACCENT)
    bullets(s, Inches(col_x), Inches(2.2), Inches(2.4), Inches(4.5),
            items, size=9, color=WHITE)

add_text(s, Inches(0), Inches(6.8), Inches(13.333), Inches(0.5),
         '特征决定上限，模型逼近上限；诊断先行，调参殿后',
         14, True, ACCENT2, align=PP_ALIGN.CENTER)

# ========== Slide 19: 结束页 ==========
s = prs.slides.add_slide(BL); set_bg(s)
for y in [3.4, 3.45, 3.5]:
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(y), Inches(2.3), Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
add_text(s, Inches(0), Inches(2.5), Inches(13.333), Inches(1.0),
         'THANKS', 72, True, WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.8), Inches(13.333), Inches(0.6),
         '感谢聆听', 24, False, GRAY_L, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.5), Inches(13.333), Inches(0.4),
         '暑期培训 · 机器学习方向 · 2026', 12, False, GRAY_M, align=PP_ALIGN.CENTER)


# ========== 保存 ==========
output_path = '/home/xavier/暑期培训/02 机器学习/机器学习学习汇报_下.pptx'
prs.save(output_path)
print(f'PPT已生成: {output_path}')
print(f'共 {len(prs.slides)} 页')
